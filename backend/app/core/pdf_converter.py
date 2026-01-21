"""
PDF Converter Module - Supports PDF to Word, Excel, and PPT
"""

import os
import tempfile
from pathlib import Path
from typing import Optional
import aiofiles
from pdf2docx import Converter
from PyPDF2 import PdfReader
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import logging

logger = logging.getLogger(__name__)


class PDFToWordConverter:
    """PDF to Word conversion handler"""

    def __init__(self, temp_dir: Optional[Path] = None):
        """
        Initialize converter

        Args:
            temp_dir: Temporary directory for file storage
        """
        if temp_dir is None:
            temp_dir = Path(tempfile.gettempdir()) / "pdf_converter"
        self.temp_dir = temp_dir
        self.temp_dir.mkdir(exist_ok=True)

    async def convert_pdf_to_word(
        self,
        pdf_path: Path,
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert PDF to Word document

        Args:
            pdf_path: Path to input PDF file
            output_filename: Optional output filename

        Returns:
            Path to converted Word document

        Raises:
            ValueError: If PDF file is invalid
            ConversionError: If conversion fails
        """
        # Validate PDF file
        if not pdf_path.exists():
            raise ValueError(f"PDF file not found: {pdf_path}")

        if not pdf_path.name.endswith('.pdf'):
            raise ValueError("File must be a PDF")

        # Check if PDF is valid
        try:
            with open(pdf_path, 'rb') as f:
                reader = PdfReader(f)
                if len(reader.pages) == 0:
                    raise ValueError("PDF has no pages")
        except Exception as e:
            raise ValueError(f"Invalid PDF file: {str(e)}")

        # Generate output filename
        if output_filename is None:
            output_filename = pdf_path.stem + '.docx'

        output_path = self.temp_dir / output_filename

        try:
            # Perform conversion
            logger.info(f"Converting {pdf_path} to {output_path}")
            cv = Converter(str(pdf_path))

            # Convert PDF to DOCX
            cv.convert(str(output_path))
            cv.close()

            logger.info(f"Conversion successful: {output_path}")

            # Verify output file was created
            if not output_path.exists():
                raise RuntimeError("Conversion failed - no output file created")

            # Check output file size
            file_size = output_path.stat().st_size
            if file_size == 0:
                raise RuntimeError("Conversion failed - empty output file")

            logger.info(f"Output file size: {file_size} bytes")

            return output_path

        except Exception as e:
            logger.error(f"Conversion error: {str(e)}")
            raise RuntimeError(f"PDF to Word conversion failed: {str(e)}") from e

    async def save_uploaded_file(self, file_content: bytes, filename: str) -> Path:
        """
        Save uploaded file to temp directory

        Args:
            file_content: File content as bytes
            filename: Original filename

        Returns:
            Path to saved file
        """
        # Sanitize filename
        safe_filename = Path(filename).name
        if not safe_filename.endswith('.pdf'):
            safe_filename += '.pdf'

        output_path = self.temp_dir / safe_filename

        # Write file
        async with aiofiles.open(output_path, 'wb') as f:
            await f.write(file_content)

        logger.info(f"Saved uploaded file: {output_path}")

        return output_path

    async def cleanup_file(self, file_path: Path) -> bool:
        """
        Delete file from temp directory

        Args:
            file_path: Path to file to delete

        Returns:
            True if successful, False otherwise
        """
        try:
            if file_path.exists():
                file_path.unlink()
                logger.info(f"Cleaned up file: {file_path}")
                return True
            return False
        except Exception as e:
            logger.error(f"Cleanup error: {str(e)}")
            return False

    def get_file_info(self, file_path: Path) -> dict:
        """
        Get file information

        Args:
            file_path: Path to file

        Returns:
            Dictionary with file information
        """
        if not file_path.exists():
            return {"error": "File not found"}

        stat = file_path.stat()

        return {
            "filename": file_path.name,
            "size": stat.st_size,
            "size_mb": round(stat.st_size / (1024 * 1024), 2),
            "extension": file_path.suffix,
        }

    def validate_pdf(self, file_path: Path) -> dict:
        """
        Validate PDF file and extract metadata

        Args:
            file_path: Path to PDF file

        Returns:
            Dictionary with validation results
        """
        try:
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                page_count = len(reader.pages)

                # Check if PDF has images (simplified check)
                has_images = False
                try:
                    for page in reader.pages:
                        if '/Resources' in page:
                            resources = page['/Resources']
                            if '/XObject' in resources:
                                has_images = True
                                break
                except Exception:
                    # If image check fails, just assume no images
                    has_images = False

                return {
                    "valid": True,
                    "pages": page_count,
                    "encrypted": reader.is_encrypted,
                    "has_images": has_images,
                }
        except Exception as e:
            return {
                "valid": False,
                "error": str(e)
            }

    async def convert_pdf_to_excel(
        self,
        pdf_path: Path,
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert PDF to Excel document
        Extracts tables from PDF and saves to Excel

        Args:
            pdf_path: Path to input PDF file
            output_filename: Optional output filename

        Returns:
            Path to converted Excel document
        """
        if not pdf_path.exists():
            raise ValueError(f"PDF file not found: {pdf_path}")

        # Generate output filename
        if output_filename is None:
            output_filename = pdf_path.stem + '.xlsx'

        output_path = self.temp_dir / output_filename

        try:
            logger.info(f"Converting {pdf_path} to Excel: {output_path}")

            # Use pdfplumber to extract tables
            all_tables = []
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    if tables:
                        for table_num, table in enumerate(tables):
                            # Convert table to DataFrame
                            df = pd.DataFrame(table[1:], columns=table[0] if table[0] else None)
                            # Add page and table info
                            df['_page'] = page_num + 1
                            df['_table'] = table_num + 1
                            all_tables.append(df)

            if all_tables:
                # Combine all tables and save to Excel
                result_df = pd.concat(all_tables, ignore_index=True)

                # Write to Excel with multiple sheets if needed
                with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                    # Group by page and write to different sheets
                    pages = result_df.groupby('_page')
                    for page_num, (page_val, page_data) in enumerate(pages, 1):
                        sheet_name = f'Page{int(page_val)}'
                        page_data = page_data.drop(['_page', '_table'], axis=1)
                        page_data.to_excel(writer, sheet_name=sheet_name, index=False, header=True)

                logger.info(f"Excel conversion successful: {output_path}")
            else:
                # No tables found, create a basic Excel with text content
                with pdfplumber.open(pdf_path) as pdf:
                    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                        for page_num, page in enumerate(pdf.pages):
                            text = page.extract_text()
                            if text:
                                df = pd.DataFrame({'Content': text.split('\n')})
                                df.to_excel(writer, sheet_name=f'Page{page_num+1}', index=False)

            # Verify output
            if not output_path.exists():
                raise RuntimeError("Excel conversion failed - no output file created")

            file_size = output_path.stat().st_size
            if file_size == 0:
                raise RuntimeError("Excel conversion failed - empty output file")

            logger.info(f"Excel conversion complete: {file_size} bytes")
            return output_path

        except Exception as e:
            logger.error(f"PDF to Excel conversion error: {str(e)}")
            raise RuntimeError(f"PDF to Excel conversion failed: {str(e)}") from e

    async def convert_pdf_to_ppt(
        self,
        pdf_path: Path,
        output_filename: Optional[str] = None
    ) -> Path:
        """
        Convert PDF to PowerPoint presentation
        Each PDF page becomes a PPT slide

        Args:
            pdf_path: Path to input PDF file
            output_filename: Optional output filename

        Returns:
            Path to converted PowerPoint presentation
        """
        if not pdf_path.exists():
            raise ValueError(f"PDF file not found: {pdf_path}")

        # Generate output filename
        if output_filename is None:
            output_filename = pdf_path.stem + '.pptx'

        output_path = self.temp_dir / output_filename

        try:
            logger.info(f"Converting {pdf_path} to PPT: {output_path}")

            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # Use pdfplumber to extract text and layout
            with pdfplumber.open(pdf_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # Extract text content
                    text = page.extract_text()

                    # Create slide
                    slide_layout = prs.slide_layouts[6]  # Blank layout
                    slide = prs.slides.add_slide(slide_layout)

                    if text:
                        # Add title with page number
                        title_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(0.5), Inches(9), Inches(0.5)
                        )
                        title_frame = title_box.text_frame
                        title_frame.text = f"Page {page_num + 1}"

                        # Add content
                        content_box = slide.shapes.add_textbox(
                            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
                        )
                        text_frame = content_box.text_frame
                        text_frame.word_wrap = True

                        # Split text into paragraphs
                        paragraphs = text.split('\n')
                        for para_text in paragraphs:
                            if para_text.strip():
                                p = text_frame.add_paragraph()
                                p.text = para_text.strip()
                                p.font.size = Pt(12)

            # Save presentation
            prs.save(output_path)

            logger.info(f"PPT conversion successful: {output_path}")

            # Verify output
            if not output_path.exists():
                raise RuntimeError("PPT conversion failed - no output file created")

            file_size = output_path.stat().st_size
            if file_size == 0:
                raise RuntimeError("PPT conversion failed - empty output file")

            logger.info(f"PPT conversion complete: {file_size} bytes")
            return output_path

        except Exception as e:
            logger.error(f"PDF to PPT conversion error: {str(e)}")
            raise RuntimeError(f"PDF to PPT conversion failed: {str(e)}") from e
